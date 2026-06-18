"""
Generate the Argus platform deck — Accenture Light theme.

Narrative arc:
  Act 0  Cover
  Act 1  Non-technical: problem -> why today's tools fall short -> the Argus difference -> what it does
  Act 2  How it works: architecture -> pipeline -> deterministic-first + AI
  Act 3  Choose your scan: modes x pipelines x approaches -> skills adaptation
  Act 4  Governance & cost
  Act 5  Use it: individual -> enterprise -> getting started
  Act 6  Reference & close: tech stack -> roadmap -> ROI -> close

Run:    .venv/bin/python scripts/generate_deck.py [output.pptx]
Output: Argus_Platform_Deck.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

# ── Accenture Light palette ────────────────────────────────────────────────────
PURPLE      = RGBColor(0xA1, 0x00, 0xFF)   # #A100FF  signature
PURPLE_BR   = RGBColor(0xC1, 0x4D, 0xFF)   # brighter purple (gradient end)
PURPLE_DK   = RGBColor(0x6A, 0x00, 0xC9)   # deep purple
PURPLE_DEEP = RGBColor(0x37, 0x00, 0x5C)   # darkest purple
PINK        = RGBColor(0xE0, 0x3C, 0xFF)   # magenta accent
INDIGO      = RGBColor(0x5A, 0x2B, 0xD9)
BLUE        = RGBColor(0x00, 0x8E, 0xF0)
TEAL        = RGBColor(0x00, 0xB5, 0xA5)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
CORAL       = RGBColor(0xFF, 0x5C, 0x39)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)

DARK        = RGBColor(0x1A, 0x1A, 0x1A)   # primary text
GRAY        = RGBColor(0x66, 0x66, 0x66)   # secondary text
GRAY_LT     = RGBColor(0x8C, 0x8C, 0x8C)
LINE        = RGBColor(0xE4, 0xE0, 0xEC)   # hairline dividers
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xFF, 0xFF, 0xFF)   # slide background
BG_TINT     = RGBColor(0xF7, 0xF4, 0xFC)   # very light purple wash
CARD        = RGBColor(0xFF, 0xFF, 0xFF)
CARD_TINT   = RGBColor(0xF4, 0xEF, 0xFB)   # light purple card

FONT        = "Arial"

# Hex strings for XML helpers
HEX_PURPLE      = "A100FF"
HEX_PURPLE_BR   = "C14DFF"
HEX_PURPLE_DK   = "6A00C9"
HEX_PURPLE_DEEP = "37005C"
HEX_PINK        = "E03CFF"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN  = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN


# ══════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _no_line(shape):
    shape.line.fill.background()


def set_gradient(shape, c1_hex: str, c2_hex: str, angle: float = 90.0):
    """Linear gradient fill. angle in degrees (0=left→right, 90=top→bottom)."""
    shape.fill.solid()
    spPr = shape._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    ang = int(angle * 60000)
    xml = (
        f'<a:gradFill {nsdecls("a")} rotWithShape="1">'
        f"<a:gsLst>"
        f'<a:gs pos="0"><a:srgbClr val="{c1_hex}"/></a:gs>'
        f'<a:gs pos="55000"><a:srgbClr val="{c2_hex}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c1_hex}"/></a:gs>'
        f"</a:gsLst>"
        f'<a:lin ang="{ang}" scaled="1"/>'
        f"</a:gradFill>"
    )
    grad = parse_xml(xml)
    if solid is not None:
        spPr.replace(solid, grad)
    else:
        spPr.append(grad)


def set_gradient_2stop(shape, c1_hex: str, c2_hex: str, angle: float = 45.0):
    shape.fill.solid()
    spPr = shape._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    ang = int(angle * 60000)
    xml = (
        f'<a:gradFill {nsdecls("a")} rotWithShape="1">'
        f"<a:gsLst>"
        f'<a:gs pos="0"><a:srgbClr val="{c1_hex}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c2_hex}"/></a:gs>'
        f"</a:gsLst>"
        f'<a:lin ang="{ang}" scaled="1"/>'
        f"</a:gradFill>"
    )
    grad = parse_xml(xml)
    if solid is not None:
        spPr.replace(solid, grad)
    else:
        spPr.append(grad)


def add_shadow(shape, alpha: int = 76, blur: int = 95000, dist: int = 38000, direction: int = 5400000):
    """Soft outer drop shadow. alpha is 0-100 (% opacity of the shadow)."""
    spPr = shape._element.spPr
    xml = (
        f'<a:effectLst {nsdecls("a")}>'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="{direction}" rotWithShape="0">'
        f'<a:srgbClr val="6A00C9"><a:alpha val="{alpha * 1000}"/></a:srgbClr>'
        f"</a:outerShdw>"
        f"</a:effectLst>"
    )
    spPr.append(parse_xml(xml))


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    else:
        _no_line(sp)
    sp.shadow.inherit = False
    return sp


def rrect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=0.08, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if fill is not None:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    else:
        _no_line(sp)
    sp.shadow.inherit = False
    if shadow:
        add_shadow(sp)
    return sp


def oval(slide, x, y, d, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is not None:
        sp.line.color.rgb = line
    else:
        _no_line(sp)
    sp.shadow.inherit = False
    return sp


def chevron(slide, x, y, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    _no_line(sp)
    sp.shadow.inherit = False
    return sp


def text(slide, x, y, w, h, s, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         italic=False, anchor=MSO_ANCHOR.TOP, font=FONT, spacing=None, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=12.5, color=DARK, gap=6, bullet_color=PURPLE,
            line_spacing=1.05, anchor=MSO_ANCHOR.TOP, bold_lead=False):
    """items: list of str. Renders dot-led bullets with breathing room."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.size = Pt(size)
        dot.font.color.rgb = bullet_color
        dot.font.name = FONT
        dot.font.bold = True
        if bold_lead and "—" in item:
            lead, rest = item.split("—", 1)
            r1 = p.add_run()
            r1.text = lead.strip() + "  "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = DARK
            r1.font.name = FONT
            r2 = p.add_run()
            r2.text = rest.strip()
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = FONT
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = FONT
    return tb


def icon_chip(slide, x, y, d, glyph, bg=PURPLE, fg=WHITE, gradient=None, glyph_size=None):
    """Colored rounded square with a centered glyph — a lightweight 'icon'."""
    chip = rrect(slide, x, y, d, d, fill=bg, radius=0.28)
    if gradient is not None:
        set_gradient_2stop(chip, gradient[0], gradient[1], 45)
    text(slide, x, y - Inches(0.01), d, d, glyph,
         size=glyph_size or int(d / Inches(1) * 19), bold=True, color=fg,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return chip


# ══════════════════════════════════════════════════════════════════════════════
#  CHROME (header / footer / dividers / backgrounds)
# ══════════════════════════════════════════════════════════════════════════════

def bg_plain(slide, color=BG):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)


def header(slide, kicker, title, on_tint=False):
    """Consistent slide header: kicker label, big title, chevron accent + hairline."""
    chevron(slide, MARGIN, Inches(0.52), Inches(0.20), Inches(0.20), PURPLE)
    text(slide, MARGIN + Inches(0.30), Inches(0.50), Inches(9), Inches(0.28),
         kicker.upper(), size=12, bold=True, color=PURPLE)
    text(slide, MARGIN, Inches(0.78), SLIDE_W - 2 * MARGIN, Inches(0.7),
         title, size=29, bold=True, color=DARK)
    rect(slide, MARGIN, Inches(1.52), SLIDE_W - 2 * MARGIN, Pt(1.4), fill=LINE)


def footer(slide, num, total):
    rect(slide, MARGIN, SLIDE_H - Inches(0.46), SLIDE_W - 2 * MARGIN, Pt(0.9), fill=LINE)
    text(slide, MARGIN, SLIDE_H - Inches(0.40), Inches(6), Inches(0.3),
         "ARGUS", size=9, bold=True, color=PURPLE)
    text(slide, MARGIN + Inches(0.62), SLIDE_H - Inches(0.40), Inches(6), Inches(0.3),
         "Adaptive AI Security Platform", size=9, color=GRAY_LT)
    text(slide, SLIDE_W - MARGIN - Inches(1.2), SLIDE_H - Inches(0.40), Inches(1.2), Inches(0.3),
         f"{num:02d} / {total:02d}", size=9, color=GRAY_LT, align=PP_ALIGN.RIGHT)


def content_slide(prs, kicker, title, tint=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_plain(slide, BG_TINT if tint else BG)
    header(slide, kicker, title)
    return slide


def _solid_alpha(shape, alpha_pct: int):
    """Make a solid-filled shape translucent. alpha_pct 0-100 (% opacity)."""
    spPr = shape._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    a = parse_xml(f'<a:alpha {nsdecls("a")} val="{alpha_pct * 1000}"/>')
    srgb.append(a)


def divider(prs, act_label, title, subtitle, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PURPLE)
    set_gradient(bg, HEX_PURPLE_DEEP, HEX_PURPLE, 35)
    for xx in [Inches(9.6), Inches(10.7), Inches(11.8)]:
        cv = chevron(slide, xx, Inches(1.4), Inches(2.2), Inches(4.6), RGBColor(0xB8, 0x55, 0xFF))
        _solid_alpha(cv, 14)
    text(slide, MARGIN, Inches(2.2), Inches(4), Inches(0.5),
         act_label.upper(), size=15, bold=True, color=RGBColor(0xE6, 0xC6, 0xFF))
    text(slide, MARGIN, Inches(2.7), Inches(11.4), Inches(1.5),
         title, size=46, bold=True, color=WHITE)
    rect(slide, MARGIN, Inches(4.18), Inches(2.4), Pt(2.6), fill=RGBColor(0xCE, 0x8B, 0xFF))
    text(slide, MARGIN, Inches(4.35), Inches(10.5), Inches(1.0),
         subtitle, size=16, color=RGBColor(0xE6, 0xCF, 0xFF), spacing=1.15)
    text(slide, SLIDE_W - MARGIN - Inches(1.2), SLIDE_H - Inches(0.55), Inches(1.2), Inches(0.3),
         f"{num:02d} / {total:02d}", size=9, color=RGBColor(0xD9, 0xBF, 0xF2), align=PP_ALIGN.RIGHT)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  REUSABLE CARD COMPONENTS
# ══════════════════════════════════════════════════════════════════════════════

def feature_card(slide, x, y, w, h, glyph, title, body, accent=PURPLE, grad=None,
                 title_size=13.5, body_size=11):
    rrect(slide, x, y, w, h, fill=CARD, radius=0.09, shadow=True)
    icon_chip(slide, x + Inches(0.22), y + Inches(0.24), Inches(0.62), glyph,
              bg=accent, gradient=grad)
    text(slide, x + Inches(0.22), y + Inches(1.02), w - Inches(0.44), Inches(0.5),
         title, size=title_size, bold=True, color=DARK)
    text(slide, x + Inches(0.22), y + Inches(1.52), w - Inches(0.44), h - Inches(1.66),
         body, size=body_size, color=GRAY, spacing=1.06)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def s_cover(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_plain(slide, BG)
    panel = rect(slide, Inches(8.0), 0, SLIDE_W - Inches(8.0), SLIDE_H, fill=PURPLE)
    set_gradient(panel, HEX_PURPLE_DEEP, HEX_PURPLE, 60)
    for xx in [Inches(10.4), Inches(11.5), Inches(12.6)]:
        cv = chevron(slide, xx, Inches(1.0), Inches(2.4), Inches(5.4), WHITE)
        _solid_alpha(cv, 13)
    chevron(slide, MARGIN, Inches(1.55), Inches(0.46), Inches(0.46), PURPLE)
    text(slide, MARGIN, Inches(2.18), Inches(7.2), Inches(1.5),
         "ARGUS", size=80, bold=True, color=DARK)
    rect(slide, MARGIN + Inches(0.05), Inches(3.5), Inches(2.6), Pt(4), fill=PURPLE)
    text(slide, MARGIN, Inches(3.72), Inches(7.0), Inches(0.6),
         "Adaptive AI Security Platform", size=23, bold=True, color=PURPLE_DK)
    text(slide, MARGIN, Inches(4.45), Inches(7.0), Inches(1.4),
         "Find vulnerabilities faster. Understand them in context. "
         "Fix them with AI — under a hard cost budget you control.",
         size=15, color=GRAY, spacing=1.2)
    pills = ["SAST", "SCA", "Secrets", "IaC", "DAST", "AI Remediation"]
    px = MARGIN
    for pgl in pills:
        wpx = Inches(0.42 + 0.115 * len(pgl))
        rrect(slide, px, Inches(5.95), wpx, Inches(0.42), fill=CARD_TINT, radius=0.5)
        text(slide, px, Inches(5.95), wpx, Inches(0.42), pgl, size=10.5, bold=True,
             color=PURPLE_DK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        px += wpx + Inches(0.12)
    text(slide, Inches(8.0), Inches(6.45), SLIDE_W - Inches(8.0), Inches(0.4),
         "Provider-agnostic  ·  Cost-aware  ·  Enterprise-ready", size=11.5,
         color=RGBColor(0xEC, 0xD6, 0xFF), align=PP_ALIGN.CENTER)
    text(slide, Inches(8.0), Inches(0.55), SLIDE_W - Inches(8.0) - Inches(0.4), Inches(0.4),
         "Accenture", size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    return slide


def s_exec_summary(prs, num, total):
    slide = content_slide(prs, "Executive Summary", "Security that keeps up — without runaway cost")
    rrect(slide, MARGIN, Inches(1.78), CONTENT_W, Inches(1.45), fill=CARD_TINT, radius=0.07)
    text(slide, MARGIN + Inches(0.35), Inches(1.78), CONTENT_W - Inches(0.7), Inches(1.45),
         "Argus unifies six industry-standard security scanners and a team of AI agents behind one API. "
         "It finds real vulnerabilities, explains them in plain language, writes the fix — and meters "
         "every AI dollar so security at scale never becomes a budget surprise.",
         size=15.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE, spacing=1.18)
    cards = [
        ("◆", "Less noise", "Deterministic scanners + AI triage cut false positives and duplicate alerts so teams fix what matters.", PURPLE, (HEX_PURPLE, HEX_PINK)),
        ("$", "Predictable cost", "Every AI call passes a budget gate. Hard per-scan and monthly limits mean no overspend, ever.", INDIGO, ("5A2BD9", "008EF0")),
        ("↻", "Fits any team", "One developer on a laptop or a whole enterprise in CI/CD — same platform, same guardrails.", TEAL, ("00B5A5", "16A34A")),
    ]
    cw = (CONTENT_W - Inches(0.7)) / 3
    for i, (g, t, b, acc, grad) in enumerate(cards):
        x = MARGIN + i * (cw + Inches(0.35))
        feature_card(slide, x, Inches(3.55), cw, Inches(3.0), g, t, b, accent=acc, grad=grad,
                     title_size=16, body_size=12)
    footer(slide, num, total)
    return slide


def s_problem(prs, num, total):
    slide = content_slide(prs, "The Problem", "Security teams are losing the race", tint=True)
    text(slide, MARGIN, Inches(1.66), CONTENT_W, Inches(0.5),
         "Code ships faster than ever. The tooling meant to secure it creates as many problems as it solves.",
         size=14, color=GRAY)
    probs = [
        ("⚑", "Alert fatigue", "Scanners emit thousands of findings per repo. Most are duplicates or false positives. Real risks drown in the noise.", CORAL),
        ("◫", "Tool sprawl", "SAST, SCA, secrets, IaC, DAST — each a separate tool, format, and dashboard. No single view of risk.", AMBER),
        ("$", "Unpredictable AI cost", "Bolting an LLM onto scanning works in a demo, then bills explode. No budget, no metering, no controls.", PINK),
        ("⚖", "No governance", "Who suppressed that finding? Which policy passed? When? Auditors ask — and nobody can answer.", INDIGO),
    ]
    cw = (CONTENT_W - Inches(0.5)) / 2
    ch = Inches(2.05)
    for i, (g, t, b, acc) in enumerate(probs):
        x = MARGIN + (i % 2) * (cw + Inches(0.5))
        y = Inches(2.3) + (i // 2) * (ch + Inches(0.32))
        rrect(slide, x, y, cw, ch, fill=CARD, radius=0.08, shadow=True)
        icon_chip(slide, x + Inches(0.26), y + Inches(0.30), Inches(0.66), g, bg=acc)
        text(slide, x + Inches(1.12), y + Inches(0.34), cw - Inches(1.3), Inches(0.5),
             t, size=15.5, bold=True, color=DARK)
        text(slide, x + Inches(1.12), y + Inches(0.86), cw - Inches(1.3), ch - Inches(1.0),
             b, size=11.5, color=GRAY, spacing=1.1)
    footer(slide, num, total)
    return slide


def s_why_shortfall(prs, num, total):
    slide = content_slide(prs, "Why Today's Solutions Fall Short", "The gap between scanning and securing")
    rows = [
        ("Raw scanners (Semgrep, Grype…)", "Fast and free, but output a flat list. No prioritization, no context, no fix. The triage burden lands on humans."),
        ("Legacy SAST suites", "Heavyweight, slow, single-language focus, and costly per seat. Findings still need manual interpretation."),
        ("Bolt-on \"AI security\" tools", "Impressive demos, opaque costs. No budget ceiling, no audit trail, and your source often leaves your boundary."),
        ("In-house glue scripts", "Stitch tools together until they break. No governance, no RBAC, no roadmap — and a maintenance tax forever."),
    ]
    y = Inches(1.82)
    rh = Inches(1.16)
    for i, (left, right) in enumerate(rows):
        yy = y + i * (rh + Inches(0.12))
        rrect(slide, MARGIN, yy, CONTENT_W, rh, fill=(BG_TINT if i % 2 else CARD), radius=0.06,
              line=LINE, line_w=Pt(0.75))
        rect(slide, MARGIN, yy, Inches(0.09), rh, fill=GRAY_LT)
        text(slide, MARGIN + Inches(0.32), yy, Inches(3.9), rh,
             left, size=13.5, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        oval(slide, MARGIN + Inches(4.35), yy + rh / 2 - Inches(0.05), Inches(0.10), PURPLE)
        text(slide, MARGIN + Inches(4.7), yy, CONTENT_W - Inches(4.9), rh,
             right, size=12, color=GRAY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    footer(slide, num, total)
    return slide


def s_difference(prs, num, total):
    slide = content_slide(prs, "The Argus Difference", "From a list of problems to a list of solved problems", tint=True)
    colw = (CONTENT_W - Inches(0.5)) / 2
    rrect(slide, MARGIN, Inches(1.85), colw, Inches(4.7), fill=CARD, radius=0.06,
          line=LINE, line_w=Pt(1))
    text(slide, MARGIN, Inches(2.05), colw, Inches(0.5), "TRADITIONAL TOOLING",
         size=13, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    trad = [
        "Thousands of raw, unranked findings",
        "Separate tool per scan type",
        "Manual triage and interpretation",
        "Fixes are the developer's problem",
        "AI cost is unbounded (if any AI at all)",
        "Suppressions live in someone's head",
        "Bring-your-own governance",
    ]
    bullets(slide, MARGIN + Inches(0.45), Inches(2.65), colw - Inches(0.8), Inches(3.7),
            trad, size=12.5, color=GRAY, gap=11, bullet_color=GRAY_LT)
    ax = MARGIN + colw + Inches(0.5)
    a = rrect(slide, ax, Inches(1.85), colw, Inches(4.7), fill=PURPLE, radius=0.06, shadow=True)
    set_gradient(a, HEX_PURPLE_DK, HEX_PURPLE, 50)
    text(slide, ax, Inches(2.05), colw, Inches(0.5), "WITH ARGUS",
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    argus = [
        "Prioritized, de-duplicated, exploit-ranked",
        "Six scanners unified behind one API",
        "AI triages, scores, and filters for you",
        "AI writes a diff-ready fix + test",
        "Hard budget gate on every AI call",
        "Suppressions are rules, versioned + audited",
        "RBAC, policies, and audit log built in",
    ]
    tb = slide.shapes.add_textbox(ax + Inches(0.45), Inches(2.65), colw - Inches(0.8), Inches(3.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(argus):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(11)
        p.line_spacing = 1.05
        c = p.add_run()
        c.text = "✓  "
        c.font.size = Pt(12.5)
        c.font.bold = True
        c.font.color.rgb = RGBColor(0xE9, 0xC9, 0xFF)
        c.font.name = FONT
        r = p.add_run()
        r.text = item
        r.font.size = Pt(12.5)
        r.font.color.rgb = WHITE
        r.font.name = FONT
    footer(slide, num, total)
    return slide


def s_what_it_does(prs, num, total):
    slide = content_slide(prs, "What Argus Does", "Four steps, fully automated")
    steps = [
        ("◎", "Find", "Six scanners sweep code, dependencies, secrets, infrastructure, and running apps — with zero AI cost.", PURPLE, ("A100FF", "E03CFF")),
        ("?", "Understand", "AI triages every finding: is it real, how exploitable, what's the blast radius? Noise is filtered out.", INDIGO, ("5A2BD9", "A100FF")),
        ("✎", "Fix", "For confirmed issues, AI drafts a precise code patch and a test — ready to review and merge.", BLUE, ("008EF0", "00B5A5")),
        ("⚖", "Govern", "Cost, policies, suppressions, and a full audit trail keep every scan compliant and affordable.", TEAL, ("00B5A5", "16A34A")),
    ]
    cw = (CONTENT_W - Inches(1.05)) / 4
    y = Inches(2.15)
    h = Inches(3.7)
    for i, (g, t, b, acc, grad) in enumerate(steps):
        x = MARGIN + i * (cw + Inches(0.35))
        rrect(slide, x, y, cw, h, fill=CARD, radius=0.09, shadow=True)
        text(slide, x + Inches(0.22), y + Inches(0.16), cw - Inches(0.4), Inches(0.4),
             f"0{i+1}", size=13, bold=True, color=LINE)
        icon_chip(slide, x + Inches(0.24), y + Inches(0.6), Inches(0.74), g,
                  bg=acc, gradient=grad, glyph_size=18)
        text(slide, x + Inches(0.24), y + Inches(1.55), cw - Inches(0.48), Inches(0.5),
             t, size=17, bold=True, color=DARK)
        text(slide, x + Inches(0.24), y + Inches(2.05), cw - Inches(0.48), h - Inches(2.2),
             b, size=11.5, color=GRAY, spacing=1.12)
        if i < 3:
            text(slide, x + cw + Inches(0.02), y + Inches(1.4), Inches(0.33), Inches(0.6),
                 "›", size=26, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    footer(slide, num, total)
    return slide


def s_architecture(prs, num, total):
    slide = content_slide(prs, "Architecture", "Layered, provider-agnostic, source stays home")
    layers = [
        ("SURFACES", "Dashboard (React)  ·  VS Code extension  ·  CI step  ·  CLI gate", PURPLE, ("A100FF", "E03CFF")),
        ("API LAYER  ·  FastAPI / Python 3.12", "REST + Server-Sent Events  ·  OpenAPI 3.1  ·  API-key auth  ·  RBAC  ·  rate limiting", PURPLE_DK, ("6A00C9", "A100FF")),
        ("ORCHESTRATOR", "Pipeline engine  ·  GovernanceGate  ·  AI agents  ·  Skills  ·  cost ledger  ·  audit log", INDIGO, ("5A2BD9", "8E2BD9")),
        ("SCANNERS  —  zero AI tokens", "Semgrep  ·  TruffleHog  ·  Grype  ·  Checkov  ·  Nuclei  ·  OWASP ZAP", TEAL, ("00B5A5", "16A34A")),
    ]
    x = MARGIN
    w = Inches(8.5)
    y = Inches(1.78)
    lh = Inches(1.04)
    for i, (title, body, acc, grad) in enumerate(layers):
        yy = y + i * (lh + Inches(0.16))
        card = rrect(slide, x, yy, w, lh, fill=acc, radius=0.08, shadow=True)
        set_gradient_2stop(card, grad[0], grad[1], 12)
        text(slide, x + Inches(0.3), yy + Inches(0.14), w - Inches(0.6), Inches(0.4),
             title, size=13.5, bold=True, color=WHITE)
        text(slide, x + Inches(0.3), yy + Inches(0.55), w - Inches(0.6), Inches(0.4),
             body, size=11, color=RGBColor(0xF0, 0xE2, 0xFF))
        if i < 3:
            text(slide, x + w / 2 - Inches(0.2), yy + lh - Inches(0.06), Inches(0.4), Inches(0.28),
                 "▼", size=11, color=GRAY_LT, align=PP_ALIGN.CENTER)
    sx = x + w + Inches(0.4)
    sw = SLIDE_W - MARGIN - sx
    rrect(slide, sx, y, sw, Inches(2.12), fill=CARD, radius=0.09, shadow=True)
    rect(slide, sx, y, Inches(0.09), Inches(2.12), fill=BLUE)
    icon_chip(slide, sx + Inches(0.24), y + Inches(0.22), Inches(0.56), "≡", bg=BLUE)
    text(slide, sx + Inches(0.24), y + Inches(0.92), sw - Inches(0.45), Inches(0.4),
         "PostgreSQL 16", size=14, bold=True, color=DARK)
    text(slide, sx + Inches(0.24), y + Inches(1.32), sw - Inches(0.45), Inches(0.75),
         "Scans, findings, fixes, policies, audit, cost ledger — SQLAlchemy 2 async + JSONB.",
         size=10.5, color=GRAY, spacing=1.05)
    fy = y + Inches(2.3)
    rrect(slide, sx, fy, sw, Inches(2.36), fill=CARD, radius=0.09, shadow=True)
    rect(slide, sx, fy, Inches(0.09), Inches(2.36), fill=PINK)
    icon_chip(slide, sx + Inches(0.24), fy + Inches(0.22), Inches(0.56), "⇄", bg=PINK)
    text(slide, sx + Inches(0.24), fy + Inches(0.92), sw - Inches(0.45), Inches(0.4),
         "finRouter Gateway", size=14, bold=True, color=DARK)
    text(slide, sx + Inches(0.24), fy + Inches(1.32), sw - Inches(0.45), Inches(1.0),
         "Provider-agnostic LLM routing, AES-256-GCM key vault, zero-retention headers, org budgets.",
         size=10.5, color=GRAY, spacing=1.05)
    footer(slide, num, total)
    return slide


def s_pipeline(prs, num, total):
    slide = content_slide(prs, "How a Scan Runs", "A six-stage pipeline — deterministic work first", tint=True)
    stages = [
        ("Ingest", "Index the target;\ndetect languages\n& frameworks", PURPLE),
        ("Scan", "Run the selected\nscanners in\nparallel", INDIGO),
        ("Triage", "Dedup, fingerprint,\nscore, drop false\npositives", BLUE),
        ("Explain", "Plain-language\nrisk + business\ncontext", TEAL),
        ("Fix", "Diff-ready patch\n+ test, human-\ngated", GREEN),
        ("Report", "SBOM, compliance,\npolicy check,\nnotify", PURPLE_DK),
    ]
    n = len(stages)
    gap = Inches(0.22)
    cw = (CONTENT_W - gap * (n - 1)) / n
    y = Inches(2.15)
    h = Inches(2.7)
    for i, (t, b, acc) in enumerate(stages):
        x = MARGIN + i * (cw + gap)
        rrect(slide, x, y, cw, h, fill=CARD, radius=0.10, shadow=True)
        rect(slide, x, y, cw, Inches(0.12), fill=acc)
        oval(slide, x + cw / 2 - Inches(0.3), y + Inches(0.32), Inches(0.6), acc)
        text(slide, x + cw / 2 - Inches(0.3), y + Inches(0.32), Inches(0.6), Inches(0.6),
             str(i + 1), size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x + Inches(0.06), y + Inches(1.04), cw - Inches(0.12), Inches(0.4),
             t, size=13.5, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.06), y + Inches(1.46), cw - Inches(0.12), Inches(1.1),
             b, size=9.8, color=GRAY, align=PP_ALIGN.CENTER, spacing=1.05)
        if i < n - 1:
            text(slide, x + cw - Inches(0.04), y + Inches(0.95), gap + Inches(0.08), Inches(0.5),
                 "›", size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    gy = Inches(5.35)
    rrect(slide, MARGIN, gy, CONTENT_W, Inches(1.05), fill=CARD_TINT, radius=0.08)
    rect(slide, MARGIN, gy, Inches(0.12), Inches(1.05), fill=PURPLE)
    icon_chip(slide, MARGIN + Inches(0.28), gy + Inches(0.24), Inches(0.56), "$", bg=PURPLE)
    text(slide, MARGIN + Inches(1.05), gy + Inches(0.16), CONTENT_W - Inches(1.3), Inches(0.4),
         "GovernanceGate wraps every AI call", size=13.5, bold=True, color=PURPLE_DEEP)
    text(slide, MARGIN + Inches(1.05), gy + Inches(0.54), CONTENT_W - Inches(1.3), Inches(0.45),
         "Before any token is spent, the gate checks the budget, picks the cheapest viable model, "
         "records the cost, and streams progress live. Over budget → remaining AI steps are skipped, scan still completes.",
         size=10.8, color=GRAY, spacing=1.05)
    footer(slide, num, total)
    return slide


def s_deterministic_ai(prs, num, total):
    slide = content_slide(prs, "Signal Quality by Design", "Cheap, deterministic work first — AI only where it pays off")
    colw = (CONTENT_W - Inches(0.5)) / 2
    rrect(slide, MARGIN, Inches(1.85), colw, Inches(4.65), fill=CARD, radius=0.07, shadow=True)
    rrect(slide, MARGIN, Inches(1.85), colw, Inches(0.72), fill=TEAL, radius=0.07)
    rect(slide, MARGIN, Inches(2.2), colw, Inches(0.37), fill=TEAL)
    text(slide, MARGIN + Inches(0.3), Inches(1.85), colw - Inches(0.6), Inches(0.72),
         "DETERMINISTIC SCANNERS", size=13.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, MARGIN + Inches(0.3), Inches(2.74), colw - Inches(0.6), Inches(0.4),
         "Zero AI tokens · runs every time", size=11, bold=True, color=TEAL)
    det = [
        "Semgrep — static analysis (SAST), 1000+ rules",
        "TruffleHog — verified secret detection",
        "Grype — dependency CVEs (SCA)",
        "Checkov — Terraform / Kubernetes / CFN (IaC)",
        "Nuclei + OWASP ZAP — running-app testing (DAST)",
        "Dedup + fingerprint so one issue is one finding",
    ]
    bullets(slide, MARGIN + Inches(0.35), Inches(3.25), colw - Inches(0.65), Inches(3.1),
            det, size=11.8, color=GRAY, gap=10, bullet_color=TEAL)
    ax = MARGIN + colw + Inches(0.5)
    rrect(slide, ax, Inches(1.85), colw, Inches(4.65), fill=CARD, radius=0.07, shadow=True)
    rrect(slide, ax, Inches(1.85), colw, Inches(0.72), fill=PURPLE, radius=0.07)
    rect(slide, ax, Inches(2.2), colw, Inches(0.37), fill=PURPLE)
    text(slide, ax + Inches(0.3), Inches(1.85), colw - Inches(0.6), Inches(0.72),
         "AI AGENTS", size=13.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, ax + Inches(0.3), Inches(2.74), colw - Inches(0.6), Inches(0.4),
         "Metered · runs only after dedup", size=11, bold=True, color=PURPLE)
    ai = [
        "Triage — confirm, score exploitability, cut FPs",
        "Explainer — plain-language risk + business impact",
        "Fix — generate a reviewed code diff and a test",
        "Pattern — connect findings into systemic risks",
        "Skill Creator — author new domain knowledge live",
        "Every call routed + budgeted via GovernanceGate",
    ]
    bullets(slide, ax + Inches(0.35), Inches(3.25), colw - Inches(0.65), Inches(3.1),
            ai, size=11.8, color=GRAY, gap=10, bullet_color=PURPLE)
    footer(slide, num, total)
    return slide


def s_choose_scan(prs, num, total):
    slide = content_slide(prs, "Choose Your Scan", "Three independent dials tune every scan", tint=True)
    text(slide, MARGIN, Inches(1.66), CONTENT_W, Inches(0.5),
         "Mode, pipeline, and approach combine freely — the same codebase can be scanned many ways for many audiences.",
         size=14, color=GRAY)
    dials = [
        ("MODE", "When & how much", "How the scan is triggered and scoped.",
         ["at_rest — full audit of a target", "batch — many targets at once", "real_time — only changed files (PR loop)"], ("A100FF", "E03CFF"), PURPLE),
        ("PIPELINE", "Which tools run", "A named recipe of scanners + agents.",
         ["full-scan · pr-check · real-time", "sca-scan · iac-scan · dast-scan", "comprehensive — everything, parallel"], ("5A2BD9", "A100FF"), INDIGO),
        ("APPROACH", "Through which lens", "The analytical mindset AI applies.",
         ["pen-test · adversary emulation", "breach & attack sim · assumed breach", "blue team · purple team"], ("00B5A5", "16A34A"), TEAL),
    ]
    cw = (CONTENT_W - Inches(0.7)) / 3
    y = Inches(2.35)
    h = Inches(3.95)
    for i, (label, sub, desc, items, grad, acc) in enumerate(dials):
        x = MARGIN + i * (cw + Inches(0.35))
        rrect(slide, x, y, cw, h, fill=CARD, radius=0.09, shadow=True)
        top = rrect(slide, x, y, cw, Inches(1.0), fill=acc, radius=0.09)
        set_gradient_2stop(top, grad[0], grad[1], 20)
        sq = rect(slide, x, y + Inches(0.5), cw, Inches(0.5), fill=acc)
        set_gradient_2stop(sq, grad[0], grad[1], 20)
        text(slide, x + Inches(0.28), y + Inches(0.16), cw - Inches(0.5), Inches(0.4),
             label, size=16, bold=True, color=WHITE)
        text(slide, x + Inches(0.28), y + Inches(0.58), cw - Inches(0.5), Inches(0.35),
             sub, size=11, color=RGBColor(0xF0, 0xE2, 0xFF))
        text(slide, x + Inches(0.28), y + Inches(1.16), cw - Inches(0.5), Inches(0.6),
             desc, size=11.5, color=DARK, bold=True, spacing=1.05)
        bullets(slide, x + Inches(0.28), y + Inches(1.92), cw - Inches(0.52), Inches(1.9),
                items, size=11, color=GRAY, gap=8, bullet_color=acc)
    footer(slide, num, total)
    return slide


def s_approaches(prs, num, total):
    slide = content_slide(prs, "Security Approaches", "Same findings, six expert perspectives")
    text(slide, MARGIN, Inches(1.64), CONTENT_W, Inches(0.5),
         "Set one parameter and AI triage + explanations re-frame entirely — serving red teams, blue teams, and "
         "control owners from a single scan.",
         size=13, color=GRAY, spacing=1.05)
    items = [
        ("Penetration Testing", "Reachability, minimal payload, blast radius, exploit chains.", CORAL, True),
        ("Adversary Emulation", "Maps each finding to MITRE ATT&CK techniques & threat actors.", PINK, False),
        ("Breach & Attack Sim", "Would WAF / SIEM / EDR catch it? Flags control gaps.", AMBER, False),
        ("Assumed Breach", "Post-compromise value: escalation, lateral movement, persistence.", INDIGO, False),
        ("Blue Team", "Detection opportunities, log sources, hardening steps.", BLUE, False),
        ("Purple Team", "Both sides: attack technique + the detection that should fire.", TEAL, False),
    ]
    cw = (CONTENT_W - Inches(0.7)) / 3
    ch = Inches(1.95)
    y0 = Inches(2.28)
    for i, (t, b, acc, is_default) in enumerate(items):
        x = MARGIN + (i % 3) * (cw + Inches(0.35))
        y = y0 + (i // 3) * (ch + Inches(0.3))
        rrect(slide, x, y, cw, ch, fill=CARD, radius=0.09, shadow=True)
        rect(slide, x, y, cw, Inches(0.11), fill=acc)
        text(slide, x + Inches(0.26), y + Inches(0.28), cw - Inches(1.3), Inches(0.5),
             t, size=14, bold=True, color=DARK)
        if is_default:
            rrect(slide, x + cw - Inches(1.18), y + Inches(0.3), Inches(0.95), Inches(0.32),
                  fill=CARD_TINT, radius=0.5)
            text(slide, x + cw - Inches(1.18), y + Inches(0.3), Inches(0.95), Inches(0.32),
                 "DEFAULT", size=8.5, bold=True, color=PURPLE_DK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x + Inches(0.26), y + Inches(0.86), cw - Inches(0.5), Inches(1.0),
             b, size=11.5, color=GRAY, spacing=1.1)
    footer(slide, num, total)
    return slide


def s_skills(prs, num, total):
    slide = content_slide(prs, "How the System Adapts", "Skills — domain knowledge that activates itself", tint=True)
    text(slide, MARGIN, Inches(1.85), Inches(5.4), Inches(0.6),
         "A skill is a focused block of security guidance injected into AI prompts at scan time.",
         size=13.5, bold=True, color=DARK, spacing=1.1)
    bullets(slide, MARGIN, Inches(2.75), Inches(5.5), Inches(3.6), [
        "Matched automatically by detected language & framework — Django code pulls Django guidance, nothing irrelevant.",
        "Ships with built-ins: Python secure coding, secrets triage, IaC hardening.",
        "Toggle active / inactive per skill — every change is audit-logged.",
        "SkillCreatorAgent writes brand-new skills at runtime for a domain you name.",
        "Generated skills are usable immediately — the platform gets smarter as you use it.",
    ], size=12, color=GRAY, gap=11, bullet_color=PURPLE)
    rx = MARGIN + Inches(5.9)
    rw = SLIDE_W - MARGIN - rx
    flow = [
        ("Scan detects: Python + Django", PURPLE),
        ("Selector matches applicable skills", INDIGO),
        ("Guidance injected into agent prompts", BLUE),
        ("Sharper triage, explanations & fixes", TEAL),
        ("Missing a domain? Creator authors one", GREEN),
    ]
    fy = Inches(1.9)
    fh = Inches(0.82)
    for i, (t, acc) in enumerate(flow):
        yy = fy + i * (fh + Inches(0.20))
        rrect(slide, rx, yy, rw, fh, fill=acc, radius=0.5, shadow=True)
        oval(slide, rx + Inches(0.16), yy + fh / 2 - Inches(0.22), Inches(0.44), WHITE)
        text(slide, rx + Inches(0.16), yy + fh / 2 - Inches(0.22), Inches(0.44), Inches(0.44),
             str(i + 1), size=14, bold=True, color=acc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, rx + Inches(0.78), yy, rw - Inches(1.0), fh,
             t, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
        if i < len(flow) - 1:
            text(slide, rx + rw / 2 - Inches(0.2), yy + fh - Inches(0.05), Inches(0.4), Inches(0.26),
                 "▼", size=10, color=PURPLE, align=PP_ALIGN.CENTER)
    footer(slide, num, total)
    return slide


def s_cost(prs, num, total):
    slide = content_slide(prs, "Cost Discipline", "AI spend is a budget, not a surprise")
    colw = (CONTENT_W - Inches(0.5)) / 2
    text(slide, MARGIN, Inches(1.8), colw, Inches(0.4), "EVERY AI CALL, EVERY TIME",
         size=12, bold=True, color=PURPLE)
    flow = [
        ("Agent requests an AI call", PURPLE),
        ("GovernanceGate checks the budget", INDIGO),
        ("Routes to the cheapest viable model tier", BLUE),
        ("Within budget → call proceeds", GREEN),
        ("Over budget → blocked, steps skipped", CORAL),
        ("Tokens + USD written to the cost ledger", PURPLE_DK),
    ]
    fy = Inches(2.25)
    fh = Inches(0.62)
    for i, (t, acc) in enumerate(flow):
        yy = fy + i * (fh + Inches(0.12))
        rrect(slide, MARGIN, yy, colw, fh, fill=CARD, radius=0.5, shadow=True)
        oval(slide, MARGIN + Inches(0.13), yy + fh / 2 - Inches(0.18), Inches(0.36), acc)
        text(slide, MARGIN + Inches(0.62), yy, colw - Inches(0.8), fh,
             t, size=11.8, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    rx = MARGIN + colw + Inches(0.5)
    tiers = [
        ("FAST tier", "Claude Haiku — cheap, high-volume work like explanations", TEAL),
        ("BALANCED tier", "Claude Sonnet — default for triage, fixes, patterns", INDIGO),
        ("TOP tier", "Claude Opus — escalation when confidence is low", PURPLE),
    ]
    text(slide, rx, Inches(1.8), colw, Inches(0.4), "TIERED MODELS — RIGHT TOOL, RIGHT PRICE",
         size=12, bold=True, color=PURPLE)
    ty = Inches(2.25)
    for i, (t, b, acc) in enumerate(tiers):
        yy = ty + i * Inches(0.92)
        rrect(slide, rx, yy, colw, Inches(0.78), fill=CARD, radius=0.10, shadow=True)
        rect(slide, rx, yy, Inches(0.10), Inches(0.78), fill=acc)
        text(slide, rx + Inches(0.28), yy + Inches(0.1), colw - Inches(0.4), Inches(0.32),
             t, size=12.5, bold=True, color=acc)
        text(slide, rx + Inches(0.28), yy + Inches(0.42), colw - Inches(0.5), Inches(0.34),
             b, size=10.5, color=GRAY)
    cy = ty + Inches(2.95)
    rrect(slide, rx, cy, colw, Inches(1.35), fill=CARD_TINT, radius=0.08)
    text(slide, rx + Inches(0.28), cy + Inches(0.16), colw - Inches(0.5), Inches(0.35),
         "Per-scan + monthly limits, soft & hard", size=12, bold=True, color=PURPLE_DEEP)
    bullets(slide, rx + Inches(0.28), cy + Inches(0.58), colw - Inches(0.55), Inches(0.7), [
        "Soft limit → live warning event at 80%",
        "Hard limit → stop spending, keep findings",
    ], size=10.5, color=GRAY, gap=4, bullet_color=PURPLE)
    footer(slide, num, total)
    return slide


def s_governance(prs, num, total):
    slide = content_slide(prs, "Governance & Security", "Built for audits, not bolted on after", tint=True)
    feats = [
        ("⚖", "Policy engine", "Thresholds on critical/high/risk score, blocked OWASP & CWE. Pass/fail in CI via exit code.", PURPLE),
        ("⊘", "Suppressions", "Rules by fingerprint, path glob, or rule ID — expirable, versioned, plus a .argusignore file.", INDIGO),
        ("◷", "Audit log", "Every privileged action records actor, before, after, timestamp. Immutable and queryable.", BLUE),
        ("⚷", "RBAC", "viewer · analyst · admin enforced per request. Orgs & workspaces for multi-tenant access.", TEAL),
        ("⛨", "DAST gate", "Live-app scanners refuse to run without a valid, non-expired target authorization.", GREEN),
        ("⧉", "Key & secret hygiene", "API keys stored as SHA-256 hashes; secrets redacted before any log, prompt, or DB write.", PINK),
    ]
    cw = (CONTENT_W - Inches(0.7)) / 3
    ch = Inches(2.15)
    y0 = Inches(1.95)
    for i, (g, t, b, acc) in enumerate(feats):
        x = MARGIN + (i % 3) * (cw + Inches(0.35))
        y = y0 + (i // 3) * (ch + Inches(0.32))
        feature_card(slide, x, y, cw, ch, g, t, b, accent=acc, title_size=13.5, body_size=11)
    footer(slide, num, total)
    return slide


def s_individual(prs, num, total):
    slide = content_slide(prs, "Individual Developer", "Stay in flow — security in the editor")
    steps = [
        ("Install & run", "docker compose up -d\nuvicorn core.api.app:app --reload"),
        ("Scan locally", "POST /api/v1/scans\n{ target_ref: \"./my-app\",\n  pipeline_config_name: \"real-time\" }"),
        ("Review findings", "GET /scans/{id}/findings?q=sql\nRanked, de-duped, explained in context."),
        ("Apply the AI fix", "POST /fixes/{id}/apply\nOpens a PR with the diff + test."),
    ]
    colw = (CONTENT_W - Inches(0.5)) / 2
    for i, (t, code) in enumerate(steps):
        yy = Inches(1.9) + i * Inches(1.18)
        rrect(slide, MARGIN, yy, colw, Inches(1.04), fill=CARD, radius=0.08, shadow=True)
        oval(slide, MARGIN + Inches(0.2), yy + Inches(0.27), Inches(0.5), PURPLE)
        text(slide, MARGIN + Inches(0.2), yy + Inches(0.27), Inches(0.5), Inches(0.5),
             str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, MARGIN + Inches(0.85), yy + Inches(0.13), colw - Inches(1.0), Inches(0.34),
             t, size=13, bold=True, color=DARK)
        text(slide, MARGIN + Inches(0.85), yy + Inches(0.47), colw - Inches(1.0), Inches(0.55),
             code, size=9.5, color=GRAY, font="Courier New", spacing=1.0)
    rx = MARGIN + colw + Inches(0.5)
    panel = rrect(slide, rx, Inches(1.9), colw, Inches(4.6), fill=PURPLE, radius=0.07, shadow=True)
    set_gradient(panel, HEX_PURPLE_DK, HEX_PURPLE, 55)
    icon_chip(slide, rx + Inches(0.35), Inches(2.2), Inches(0.7), "{ }", bg=WHITE, fg=PURPLE_DK, glyph_size=15)
    text(slide, rx + Inches(0.35), Inches(3.05), colw - Inches(0.7), Inches(0.5),
         "VS Code Extension", size=20, bold=True, color=WHITE)
    perks = [
        "Inline finding highlights as you type",
        "One-click AI fix application",
        "Real-time diff-only scanning",
        "Live cost budget in the status bar",
        "Zero context-switching — never leave the editor",
    ]
    tb = slide.shapes.add_textbox(rx + Inches(0.35), Inches(3.7), colw - Inches(0.7), Inches(2.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, p_ in enumerate(perks):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(11)
        c = p.add_run(); c.text = "✓  "
        c.font.size = Pt(12.5); c.font.bold = True; c.font.name = FONT
        c.font.color.rgb = RGBColor(0xEA, 0xCB, 0xFF)
        r = p.add_run(); r.text = p_
        r.font.size = Pt(12.5); r.font.color.rgb = WHITE; r.font.name = FONT
    footer(slide, num, total)
    return slide


def s_enterprise(prs, num, total):
    slide = content_slide(prs, "Enterprise Deployment", "Org-wide coverage with guardrails", tint=True)
    scen = [
        ("⛓", "CI/CD gating", "Webhooks scan every PR. ci-gate blocks merges that breach policy. Nightly cron scans the whole estate.", PURPLE),
        ("◫", "Security operations", "Orgs & workspaces per team. Analysts triage, admins configure. Bulk-assign, suppress, and dismiss at scale.", INDIGO),
        ("✓", "Compliance & reporting", "OWASP / CWE reports, CycloneDX SBOM per scan, CSV export, and persisted policy evaluations for auditors.", TEAL),
        ("◎", "Observability & SRE", "Prometheus metrics, OpenTelemetry tracing, PagerDuty incidents, Jira tickets, Slack alerts — wired in.", BLUE),
    ]
    cw = (CONTENT_W - Inches(0.5)) / 2
    ch = Inches(2.15)
    y0 = Inches(1.95)
    for i, (g, t, b, acc) in enumerate(scen):
        x = MARGIN + (i % 2) * (cw + Inches(0.5))
        y = y0 + (i // 2) * (ch + Inches(0.32))
        rrect(slide, x, y, cw, ch, fill=CARD, radius=0.08, shadow=True)
        icon_chip(slide, x + Inches(0.28), y + Inches(0.3), Inches(0.66), g, bg=acc)
        text(slide, x + Inches(1.14), y + Inches(0.34), cw - Inches(1.35), Inches(0.5),
             t, size=15, bold=True, color=DARK)
        text(slide, x + Inches(1.14), y + Inches(0.86), cw - Inches(1.35), ch - Inches(1.0),
             b, size=11.5, color=GRAY, spacing=1.12)
    footer(slide, num, total)
    return slide


def s_getting_started(prs, num, total):
    slide = content_slide(prs, "Getting Started", "From clone to first finding in five minutes")
    steps = [
        ("1", "Prerequisites", "Python 3.12+ · Docker · uv · an ANTHROPIC_API_KEY", PURPLE),
        ("2", "Install", "git clone … && cd argus  ·  uv venv && uv pip install -e \".[dev]\"", INDIGO),
        ("3", "Start infrastructure", "docker compose up -d   (PostgreSQL + gateway)  ·  alembic upgrade head", BLUE),
        ("4", "Run the API", "uvicorn core.api.app:app --reload --port 8000", TEAL),
        ("5", "First scan", "curl -X POST localhost:8000/api/v1/scans -d '{\"target_ref\":\".\",\"pipeline_config_name\":\"full-scan\"}'", GREEN),
        ("6", "Explore", "Open localhost:8000/docs  →  interactive Swagger UI", PURPLE_DK),
    ]
    y = Inches(1.9)
    rh = Inches(0.74)
    for i, (n, t, code, acc) in enumerate(steps):
        yy = y + i * (rh + Inches(0.10))
        rrect(slide, MARGIN, yy, CONTENT_W, rh, fill=(BG_TINT if i % 2 else CARD), radius=0.06,
              line=LINE, line_w=Pt(0.75))
        oval(slide, MARGIN + Inches(0.18), yy + rh / 2 - Inches(0.24), Inches(0.48), acc)
        text(slide, MARGIN + Inches(0.18), yy + rh / 2 - Inches(0.24), Inches(0.48), Inches(0.48),
             n, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, MARGIN + Inches(0.85), yy, Inches(2.7), rh,
             t, size=13, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, MARGIN + Inches(3.7), yy, CONTENT_W - Inches(3.9), rh,
             code, size=10, color=GRAY, font="Courier New", anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
    footer(slide, num, total)
    return slide


def s_tech_stack(prs, num, total):
    slide = content_slide(prs, "Technology Stack", "Modern, typed, async, observable", tint=True)
    cats = [
        ("Backend", ["Python 3.12 · FastAPI", "Pydantic v2 (strict)", "SQLAlchemy 2 async · asyncpg", "Alembic migrations", "structlog"], PURPLE),
        ("AI / LLM", ["Anthropic Claude (3 tiers)", "Provider-agnostic routing", "GovernanceGate budgets", "Per-scan token ledger", "Batch mode savings"], INDIGO),
        ("Scanners", ["Semgrep (SAST)", "TruffleHog (secrets)", "Grype (SCA)", "Checkov (IaC)", "Nuclei + ZAP (DAST)"], TEAL),
        ("Data", ["PostgreSQL 16", "JSONB for findings", "Cursor pagination", "Full-text search", "CycloneDX SBOM"], BLUE),
        ("Observability", ["Prometheus metrics", "OpenTelemetry / OTLP", "Server-Sent Events", "Structured audit log", "structlog JSON"], GREEN),
        ("Security", ["API-key auth (SHA-256)", "RBAC roles", "HMAC webhook verify", "Rate limiting", "Secret redaction"], PINK),
    ]
    cw = (CONTENT_W - Inches(0.7)) / 3
    ch = Inches(2.25)
    y0 = Inches(1.85)
    for i, (cat, items, acc) in enumerate(cats):
        x = MARGIN + (i % 3) * (cw + Inches(0.35))
        y = y0 + (i // 3) * (ch + Inches(0.3))
        rrect(slide, x, y, cw, ch, fill=CARD, radius=0.08, shadow=True)
        rrect(slide, x, y, cw, Inches(0.52), fill=acc, radius=0.08)
        rect(slide, x, y + Inches(0.26), cw, Inches(0.26), fill=acc)
        text(slide, x + Inches(0.24), y, cw - Inches(0.4), Inches(0.52),
             cat, size=13.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        bullets(slide, x + Inches(0.26), y + Inches(0.68), cw - Inches(0.5), Inches(1.5),
                items, size=10.8, color=GRAY, gap=5, bullet_color=acc)
    footer(slide, num, total)
    return slide


def s_roadmap(prs, num, total):
    slide = content_slide(prs, "Delivery Roadmap", "Fifteen phases shipped — 433 tests green")
    phases = [
        ("1–3", "Foundation", "SAST + secrets, triage, explainer, fixes, cost ledger, live trace, VS Code extension"),
        ("4–5", "Coverage", "SCA (Grype), IaC (Checkov), batch API, skills system, pattern & skill-creator agents"),
        ("6–7", "DAST & audit", "Nuclei + ZAP with authorization gate, audit log, config API, evaluation harness"),
        ("8–9", "Integrate", "CycloneDX SBOM, scan diff, GitHub/GitLab webhooks, API-key auth, Prometheus metrics"),
        ("10–11", "Govern", "Suppressions + .argusignore, cron scheduler, compliance reports, policy engine, CI gate"),
        ("12–13", "Enterprise", "Orgs / workspaces, RBAC roles, Jira + PagerDuty + Slack integrations"),
        ("14–15", "Production", "Trend analytics, MTTR, CSV export, cursor pagination, full-text search, OpenTelemetry"),
    ]
    y = Inches(1.8)
    rh = Inches(0.62)
    for i, (ph, title, desc) in enumerate(phases):
        yy = y + i * (rh + Inches(0.115))
        badge = rrect(slide, MARGIN, yy, Inches(1.05), rh, fill=PURPLE, radius=0.18)
        set_gradient_2stop(badge, HEX_PURPLE_DK, HEX_PURPLE, 30)
        text(slide, MARGIN, yy, Inches(1.05), rh, ph, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rrect(slide, MARGIN + Inches(1.2), yy, Inches(2.3), rh, fill=CARD_TINT, radius=0.1)
        text(slide, MARGIN + Inches(1.4), yy, Inches(2.0), rh, title, size=13, bold=True,
             color=PURPLE_DEEP, anchor=MSO_ANCHOR.MIDDLE)
        rrect(slide, MARGIN + Inches(3.65), yy, CONTENT_W - Inches(4.45), rh, fill=CARD, radius=0.08,
              line=LINE, line_w=Pt(0.75))
        text(slide, MARGIN + Inches(3.9), yy, CONTENT_W - Inches(5.4), rh, desc, size=10.8,
             color=GRAY, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
        oval(slide, SLIDE_W - MARGIN - Inches(0.5), yy + rh / 2 - Inches(0.18), Inches(0.36), GREEN)
        text(slide, SLIDE_W - MARGIN - Inches(0.5), yy + rh / 2 - Inches(0.18), Inches(0.36), Inches(0.36),
             "✓", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, num, total)
    return slide


def s_roi(prs, num, total):
    slide = content_slide(prs, "Value & ROI", "Better signal, lower cost, full traceability", tint=True)
    stats = [
        ("≈70%", "less manual triage", "dedup + AI filtering remove duplicate & false-positive noise", PURPLE),
        ("~$0.50", "AI cost per full scan", "tiered models + hard budget gate keep spend predictable", INDIGO),
        ("< 5 min", "to a diff-ready fix", "for a confirmed high-severity finding, with a test", TEAL),
        ("100%", "of actions audited", "every suppression, policy, and fix is traceable", BLUE),
    ]
    cw = (CONTENT_W - Inches(1.05)) / 4
    for i, (v, l, sub, acc) in enumerate(stats):
        x = MARGIN + i * (cw + Inches(0.35))
        y = Inches(1.95)
        h = Inches(2.3)
        rrect(slide, x, y, cw, h, fill=CARD, radius=0.10, shadow=True)
        rect(slide, x, y, cw, Inches(0.11), fill=acc)
        text(slide, x + Inches(0.1), y + Inches(0.3), cw - Inches(0.2), Inches(0.8),
             v, size=34, bold=True, color=acc, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.12), y + Inches(1.15), cw - Inches(0.24), Inches(0.4),
             l, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.12), y + Inches(1.55), cw - Inches(0.24), Inches(0.7),
             sub, size=9.8, color=GRAY, align=PP_ALIGN.CENTER, spacing=1.05)
    band = rrect(slide, MARGIN, Inches(4.65), CONTENT_W, Inches(1.55), fill=PURPLE, radius=0.07, shadow=True)
    set_gradient(band, HEX_PURPLE_DK, HEX_PURPLE, 50)
    text(slide, MARGIN + Inches(0.5), Inches(4.65), CONTENT_W - Inches(1.0), Inches(1.55),
         "One platform replaces a fragmented toolchain — unifying detection, AI remediation, cost control, "
         "and governance so security scales with the business instead of fighting it.",
         size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, spacing=1.15)
    footer(slide, num, total)
    return slide


def s_closing(prs, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PURPLE)
    set_gradient(bg, HEX_PURPLE_DEEP, HEX_PURPLE, 40)
    for xx in [Inches(9.7), Inches(10.8), Inches(11.9)]:
        cv = chevron(slide, xx, Inches(0.0), Inches(2.6), SLIDE_H, WHITE)
        _solid_alpha(cv, 10)
    chevron(slide, MARGIN, Inches(2.05), Inches(0.5), Inches(0.5), RGBColor(0xD9, 0x9A, 0xFF))
    text(slide, MARGIN, Inches(2.6), Inches(10), Inches(1.4),
         "ARGUS", size=78, bold=True, color=WHITE)
    rect(slide, MARGIN + Inches(0.05), Inches(4.0), Inches(2.6), Pt(4), fill=RGBColor(0xCE, 0x8B, 0xFF))
    text(slide, MARGIN, Inches(4.2), Inches(11), Inches(0.6),
         "Find faster. Understand deeper. Fix smarter. Spend less.",
         size=20, bold=True, color=WHITE)
    text(slide, MARGIN, Inches(5.05), Inches(11), Inches(0.5),
         "Provider-agnostic  ·  cost-aware  ·  enterprise-ready security, powered by AI.",
         size=14, color=RGBColor(0xE6, 0xCF, 0xFF))
    text(slide, MARGIN, Inches(6.4), Inches(8), Inches(0.4),
         "github.com/ahujrajat/Argus", size=13, bold=True, color=RGBColor(0xEA, 0xD4, 0xFF))
    text(slide, SLIDE_W - MARGIN - Inches(3.2), Inches(6.4), Inches(3.2), Inches(0.4),
         "Built by Rajat Ahuja · Accenture", size=12, color=RGBColor(0xD9, 0xBF, 0xF2),
         align=PP_ALIGN.RIGHT)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  ASSEMBLY
# ══════════════════════════════════════════════════════════════════════════════

def build(output: str = "Argus_Platform_Deck.pptx") -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    plan = [
        ("cover", s_cover),
        ("content", s_exec_summary),
        ("content", s_problem),
        ("content", s_why_shortfall),
        ("content", s_difference),
        ("content", s_what_it_does),
        ("divider", ("ACT II", "How It Works", "Architecture, the scan pipeline, and why the signal is clean.")),
        ("content", s_architecture),
        ("content", s_pipeline),
        ("content", s_deterministic_ai),
        ("divider", ("ACT III", "Choose Your Scan", "Mode, pipeline, approach — and a system that adapts with skills.")),
        ("content", s_choose_scan),
        ("content", s_approaches),
        ("content", s_skills),
        ("divider", ("ACT IV", "Govern & Control", "Predictable cost and audit-ready governance, by design.")),
        ("content", s_cost),
        ("content", s_governance),
        ("divider", ("ACT V", "Put It to Work", "How an individual developer and a whole enterprise deploy Argus.")),
        ("content", s_individual),
        ("content", s_enterprise),
        ("content", s_getting_started),
        ("divider", ("ACT VI", "Under the Hood", "Stack, delivery roadmap, and the business case.")),
        ("content", s_tech_stack),
        ("content", s_roadmap),
        ("content", s_roi),
        ("closing", s_closing),
    ]
    total = len(plan)
    for i, (kind, payload) in enumerate(plan):
        n = i + 1
        if kind == "divider":
            act, title, sub = payload
            divider(prs, act, title, sub, n, total)
        else:
            payload(prs, n, total)

    prs.save(output)
    print(f"Saved: {output}  ({total} slides)")


if __name__ == "__main__":
    import sys
    build(sys.argv[1] if len(sys.argv) > 1 else "Argus_Platform_Deck.pptx")
